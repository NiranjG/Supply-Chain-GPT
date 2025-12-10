"""
Document parsers for various file formats
"""

import hashlib
import re
from pathlib import Path
from typing import Dict, Any, Optional, Tuple
from io import BytesIO

import pandas as pd
from pdfminer.high_level import extract_text as pdf_extract_text
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from docx import Document as DocxDocument
from pptx import Presentation


class DocumentParser:
    """Parse documents from various formats to plain text"""

    def __init__(self):
        self.supported_formats = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_xlsx,
            ".xls": self._parse_xlsx,
            ".csv": self._parse_csv,
            ".txt": self._parse_txt,
            ".md": self._parse_txt,
        }

    def parse(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Parse a document and return text content with metadata

        Args:
            file_path: Path to the document

        Returns:
            Tuple of (text_content, metadata)
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()

        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {suffix}")

        parser_func = self.supported_formats[suffix]
        text, extra_meta = parser_func(file_path)

        # Compute checksum
        with open(file_path, "rb") as f:
            checksum = hashlib.sha256(f.read()).hexdigest()

        metadata = {
            "source_uri": str(file_path),
            "doc_title": file_path.stem,
            "file_type": suffix,
            "checksum": checksum,
            "parser": parser_func.__name__,
            **extra_meta
        }

        # Clean text
        text = self._clean_text(text)

        return text, metadata

    def parse_bytes(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Parse document from bytes"""
        suffix = Path(filename).suffix.lower()

        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {suffix}")

        # Create a BytesIO object
        buffer = BytesIO(content)

        if suffix == ".pdf":
            text = pdf_extract_text(buffer)
            extra_meta = {}
        elif suffix in [".docx", ".doc"]:
            doc = DocxDocument(buffer)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            extra_meta = {}
        elif suffix == ".pptx":
            prs = Presentation(buffer)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(f"[Slide {slide_num}] {shape.text}")
            text = "\n\n".join(text_parts)
            extra_meta = {"slide_count": len(prs.slides)}
        elif suffix in [".xlsx", ".xls"]:
            df = pd.read_excel(buffer)
            text = self._dataframe_to_markdown(df)
            extra_meta = {"row_count": len(df), "columns": list(df.columns)}
        elif suffix == ".csv":
            df = pd.read_csv(buffer)
            text = self._dataframe_to_markdown(df)
            extra_meta = {"row_count": len(df), "columns": list(df.columns)}
        else:
            text = content.decode("utf-8", errors="ignore")
            extra_meta = {}

        checksum = hashlib.sha256(content).hexdigest()

        metadata = {
            "source_uri": filename,
            "doc_title": Path(filename).stem,
            "file_type": suffix,
            "checksum": checksum,
            "parser": "parse_bytes",
            **extra_meta
        }

        return self._clean_text(text), metadata

    def _parse_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse PDF document"""
        text = pdf_extract_text(str(file_path))

        # Try to get PDF metadata
        extra_meta = {}
        try:
            with open(file_path, "rb") as f:
                parser = PDFParser(f)
                doc = PDFDocument(parser)
                if doc.info:
                    info = doc.info[0] if doc.info else {}
                    if b"Title" in info:
                        extra_meta["pdf_title"] = info[b"Title"].decode("utf-8", errors="ignore")
                    if b"Author" in info:
                        extra_meta["author"] = info[b"Author"].decode("utf-8", errors="ignore")
        except Exception:
            pass

        return text, extra_meta

    def _parse_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse DOCX document"""
        doc = DocxDocument(str(file_path))

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Add heading markers
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        level = int(level)
                        prefix = "#" * level + " "
                    except ValueError:
                        prefix = "## "
                    paragraphs.append(f"{prefix}{para.text}")
                else:
                    paragraphs.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_text = self._table_to_markdown(table)
            paragraphs.append(table_text)

        text = "\n\n".join(paragraphs)

        extra_meta = {}
        try:
            core_props = doc.core_properties
            if core_props.title:
                extra_meta["doc_title"] = core_props.title
            if core_props.author:
                extra_meta["author"] = core_props.author
        except Exception:
            pass

        return text, extra_meta

    def _table_to_markdown(self, table) -> str:
        """Convert a docx table to markdown format"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if len(rows) > 1:
            # Add header separator
            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            rows.insert(1, header_sep)

        return "\n".join(rows)

    def _parse_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse PPTX presentation"""
        prs = Presentation(str(file_path))

        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_parts.append("\n".join(slide_text))

        text = "\n\n".join(text_parts)
        extra_meta = {"slide_count": len(prs.slides)}

        return text, extra_meta

    def _parse_xlsx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse Excel spreadsheet"""
        xl = pd.ExcelFile(file_path)

        text_parts = []
        total_rows = 0
        all_columns = []

        for sheet_name in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=sheet_name)
            total_rows += len(df)
            all_columns.extend(list(df.columns))

            text_parts.append(f"## Sheet: {sheet_name}")
            text_parts.append(self._dataframe_to_markdown(df))

        text = "\n\n".join(text_parts)
        extra_meta = {
            "sheet_count": len(xl.sheet_names),
            "row_count": total_rows,
            "columns": list(set(all_columns))
        }

        return text, extra_meta

    def _parse_csv(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse CSV file"""
        df = pd.read_csv(file_path)
        text = self._dataframe_to_markdown(df)

        extra_meta = {
            "row_count": len(df),
            "columns": list(df.columns)
        }

        return text, extra_meta

    def _parse_txt(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Parse plain text file"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return text, {}

    def _dataframe_to_markdown(self, df: pd.DataFrame, max_rows: int = 100) -> str:
        """Convert DataFrame to markdown table"""
        if len(df) > max_rows:
            df = df.head(max_rows)
            truncated = True
        else:
            truncated = False

        # Create header
        headers = "| " + " | ".join(str(col) for col in df.columns) + " |"
        separator = "| " + " | ".join(["---"] * len(df.columns)) + " |"

        # Create rows
        rows = []
        for _, row in df.iterrows():
            row_str = "| " + " | ".join(str(val) for val in row.values) + " |"
            rows.append(row_str)

        result = "\n".join([headers, separator] + rows)

        if truncated:
            result += f"\n\n*[Table truncated. Showing {max_rows} of {len(df)} rows]*"

        return result

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove excessive whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" {2,}", " ", text)

        # Remove control characters except newlines and tabs
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

        return text.strip()

    def extract_entities(self, text: str) -> Dict[str, Any]:
        """Extract business entities from text (SKUs, warehouses, etc.)"""
        entities = {
            "sku_ids": [],
            "warehouse_ids": [],
            "supplier_ids": [],
            "dates": [],
        }

        # SKU patterns (e.g., SKU-12345, SKU12345)
        sku_pattern = r"\b(?:SKU[-_]?\d{4,}|[A-Z]{2,3}[-_]?\d{4,})\b"
        entities["sku_ids"] = list(set(re.findall(sku_pattern, text, re.IGNORECASE)))

        # Warehouse patterns (e.g., WH-001, Warehouse-A)
        wh_pattern = r"\b(?:WH[-_]?\d{2,}|Warehouse[-_]?[A-Z0-9]+)\b"
        entities["warehouse_ids"] = list(set(re.findall(wh_pattern, text, re.IGNORECASE)))

        # Supplier patterns (e.g., SUP-001, Supplier-ABC)
        sup_pattern = r"\b(?:SUP[-_]?\d{2,}|Supplier[-_]?[A-Z0-9]+)\b"
        entities["supplier_ids"] = list(set(re.findall(sup_pattern, text, re.IGNORECASE)))

        # Date patterns
        date_pattern = r"\b\d{4}[-/]\d{2}[-/]\d{2}\b|\b\d{2}[-/]\d{2}[-/]\d{4}\b"
        entities["dates"] = list(set(re.findall(date_pattern, text)))

        return entities
